#!/usr/bin/env python3
"""Check a filled ritual deck before handing it over.

This exists because the failure modes are boring and repetitive: a token left
unfilled, the READ ME slide still in there, a slide that quietly grew into a
paragraph, a slide with no speaker notes. Catching them costs a second here and
saves the consultant a re-read.

  python audit_deck.py deck.pptx --ritual trex

Exit code is 1 when something needs a human decision, 0 when the deck is clean.
Leftover [brackets] are reported but are not automatically a failure: a
deliberate [your number here] is the honest answer when a figure is missing.
They are listed so you can say them out loud in the handover message.
"""

import argparse
import re
import sys

from pptx import Presentation

SLIDE_CEILING = {
    "trex": 7,
    "tlab": 6,
    "thats-a-banger": 7,
    "passion-fruit": 6,
    "back-2-back": 7,
    "ask-me-anything": 5,   # title + 4
    "lets-workshop": 8,
    "change-my-mind": 5,
}

WORD_LIMIT = 25          # one idea per slide
CHROME = re.compile(
    r"^(WWW\.THIGA\.CO|TRIBE DESIGN ·|\d+$|\d+ (min|sec)$)", re.IGNORECASE)
# Dashed zones where the consultant drops a visual. Their label is scaffolding,
# not something the room reads, so it must not count against the word budget.
DROP_ZONE = re.compile(
    r"^(IMAGE\b|PHOTO\b|GUEST PHOTO|DROP YOUR ARTEFACTS|SCREENSHOT OF|"
    r"ONE ZOOMED CROP|YOUR BEST PHOTOS|one archive image|diagram, photo)",
    re.IGNORECASE)
TIMING = re.compile(r"(⏱|\b\d+\s*(min|sec|minutes|seconds)\b)", re.IGNORECASE)


def body_words(slide):
    """Words the audience actually reads.

    Skips footers, page numbers, section headers and visual drop zones, and
    ignores unfilled [tokens] so the count reflects real prose rather than
    template scaffolding.
    """
    words = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or CHROME.match(text) or DROP_ZONE.match(text):
            continue
        words += len(re.sub(r"\[[^\[\]]*\]", "", text).split())
    return words


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--ritual", help=", ".join(sorted(SLIDE_CEILING)))
    args = ap.parse_args()

    prs = Presentation(args.deck)
    slides = list(prs.slides)
    problems, notices = [], []

    ceiling = SLIDE_CEILING.get(args.ritual)
    print(f"slides: {len(slides)}" + (f" (ceiling {ceiling})" if ceiling else ""))
    if ceiling and len(slides) > ceiling:
        problems.append(
            f"{len(slides)} slides for a ceiling of {ceiling}. The ceiling is a "
            f"ceiling: cut, don't negotiate.")

    for i, slide in enumerate(slides, start=1):
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        for shape in slide.shapes:
            if shape.has_table:
                texts += [c.text for r in shape.table.rows for c in r.cells]
        joined = "\n".join(texts)

        if "READ ME FIRST" in joined or "DELETE THIS SLIDE" in joined:
            problems.append(f"slide {i}: the READ ME slide is still in the deck")
        for line in joined.split("\n"):
            if re.match(r"^\s*[·•\-\s]*TIP\b", line, re.IGNORECASE):
                problems.append(f"slide {i}: a TIP box is still there")
                break

        left = sorted(set(re.findall(r"\[[^\[\]]+\]", joined)))
        if left:
            notices.append(f"slide {i}: still bracketed → {', '.join(left)}")

        w = body_words(slide)
        if w > WORD_LIMIT:
            notices.append(
                f"slide {i}: ~{w} words on screen. Over {WORD_LIMIT} usually "
                f"means a paragraph that belongs in the speaker notes.")

        notes = slide.notes_slide.notes_text_frame.text.strip() \
            if slide.has_notes_slide else ""
        if not notes:
            problems.append(f"slide {i}: no speaker notes")
        elif not TIMING.search(notes):
            notices.append(f"slide {i}: speaker notes carry no timing cue")

    print()
    for n in notices:
        print(f"  note   {n}")
    for p in problems:
        print(f"  FIX    {p}")
    if not problems and not notices:
        print("  clean")
    print()

    if notices and not problems:
        print("Nothing blocking. Repeat the bracketed items to the consultant "
              "in the handover message — those are theirs to fill.")
    sys.exit(1 if problems else 0)


if __name__ == "__main__":
    main()
