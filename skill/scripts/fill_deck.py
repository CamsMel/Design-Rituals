#!/usr/bin/env python3
"""Fill a Tribe Design ritual template and produce a presentable .pptx.

The templates are Google Slides exports: every text box is an autoshape, there
are no placeholders, and the brand formatting lives in the runs. So we never
rebuild a slide — we swap the text inside existing runs and keep everything else
untouched. That is the whole point: the deck stays on brand because we did not
redesign it.

Three jobs, in this order:
  1. delete the "READ ME FIRST" slide (it is speaker scaffolding, not content)
  2. replace the [bracketed] tokens with the consultant's own words
  3. delete the TIP boxes and write the speaker notes

Usage
-----
  # see what is fillable before writing anything
  python fill_deck.py --template trex --list

  # fill
  python fill_deck.py --template trex --data content.json --out deck.pptx

content.json
------------
{
  "slides": [
    {
      "n": 1,
      "replace": {
        "[DATE]": "17 September 2026",
        "[The story you want to tell, in one line]": "..."
      },
      "notes": "What to say out loud.\n2 min"
    },
    {
      "n": 5,
      "replace": {"[Move 1, and the artefact that carried it]": "...",
                  "[Move 3, only if it deserves the airtime]": null},
      "tables": [{"index": 0, "rows": [[null, "Yours"], ["Sector & size", "..."]]}]
    }
  ]
}

Slide numbers `n` are 1-based **in the final deck**, i.e. after the READ ME
slide is gone: slide 1 is always the cover.

Values:
  "text"        -> replaces the token, keeps the run's formatting
  null          -> deletes that whole paragraph (use it to drop an unused bullet)
  ["a", "b"]    -> for a token that appears several times: fills the occurrences
                   in document order, null entries delete their paragraph

Tables: `rows` mirrors the table grid. null in a cell means "leave it alone".
Extra rows beyond the table's size are ignored, so check with --list first.
"""

import argparse
import copy
import json
import os
import re
import sys

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DIR = os.path.normpath(os.path.join(HERE, "..", "assets", "templates"))

RITUALS = {
    "trex": "trex.pptx",
    "tlab": "tlab.pptx",
    "thats-a-banger": "thats-a-banger.pptx",
    "passion-fruit": "passion-fruit.pptx",
    "back-2-back": "back-2-back.pptx",
    "ask-me-anything": "ask-me-anything.pptx",
    "lets-workshop": "lets-workshop.pptx",
    "change-my-mind": "change-my-mind.pptx",
}

TIP_RE = re.compile(r"^\s*[·•\-\s]*TIP\b", re.IGNORECASE)


# --------------------------------------------------------------------------
# low-level text handling
# --------------------------------------------------------------------------
def set_paragraph_text(paragraph, text):
    """Write `text` into the paragraph, keeping the first run's formatting."""
    runs = paragraph.runs
    if not runs:
        # no run to inherit from; add one so the text at least lands
        paragraph.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r.text = ""


def delete_paragraph(paragraph):
    el = paragraph._p
    el.getparent().remove(el)


def iter_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape, shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield shape, cell.text_frame


def replace_token(slide, token, value):
    """Replace every occurrence of `token` on the slide, in document order.

    Counting occurrences rather than paragraphs matters: some templates put the
    same token twice on one line (`[Your name · Your client] x [Your name ·
    Your client]` on the Back 2 Back cover), and both speakers are not the same
    person. Returns the number of occurrences handled.
    """
    is_list = isinstance(value, list)
    values = value if is_list else [value]
    KEEP = object()
    idx = 0
    targets = []
    for _, tf in iter_text_frames(slide):
        for para in list(tf.paragraphs):
            if token in para.text:
                targets.append(para)
    for para in targets:
        text = para.text
        chunk = []
        for _ in range(text.count(token)):
            if idx < len(values):
                chunk.append(values[idx])
            else:
                # A scalar applies everywhere. A list that runs short leaves the
                # rest bracketed on purpose — a visible gap the audit will catch
                # beats silently duplicating the last answer.
                chunk.append(values[-1] if not is_list else KEEP)
            idx += 1
        if any(c is None for c in chunk):
            delete_paragraph(para)
            continue
        for repl in chunk:
            if repl is KEEP:
                break
            text = text.replace(token, str(repl), 1)
        set_paragraph_text(para, text)
    return idx


def fill_table(shape, rows):
    table = shape.table
    for r_i, row in enumerate(rows):
        if r_i >= len(table.rows):
            break
        for c_i, val in enumerate(row):
            if val is None or c_i >= len(table.columns):
                continue
            cell = table.cell(r_i, c_i)
            paras = cell.text_frame.paragraphs
            lines = str(val).split("\n")
            set_paragraph_text(paras[0], lines[0])
            for extra in paras[1:]:
                delete_paragraph(extra)
            for line in lines[1:]:
                new_p = copy.deepcopy(paras[0]._p)
                paras[0]._p.getparent().append(new_p)
                set_paragraph_text(cell.text_frame.paragraphs[-1], line)


# --------------------------------------------------------------------------
# cleanup
# --------------------------------------------------------------------------
def drop_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


EMU_IN = 914400


def backs_onto(bar, tip, pad=0.35):
    """True when `bar` is plausibly the coloured band drawn behind `tip`.

    The band is always a little wider than the text box it frames and shares
    most of its vertical extent, so test containment rather than equality —
    the two shapes are offset by design and never line up exactly.
    """
    vals = [bar.left, bar.top, bar.width, bar.height,
            tip.left, tip.top, tip.width, tip.height]
    if any(v is None for v in vals):
        return False
    pad_emu = int(pad * EMU_IN)
    if bar.left > tip.left + pad_emu:
        return False
    if bar.left + bar.width < tip.left + tip.width - pad_emu:
        return False
    top = max(bar.top, tip.top)
    bottom = min(bar.top + bar.height, tip.top + tip.height)
    return (bottom - top) > 0.5 * tip.height


def strip_tips(slide):
    """Remove TIP callouts, text and the coloured bar drawn behind them.

    The bar is a separate borderless autoshape sitting at almost the same
    coordinates; leaving it behind puts an empty stripe at the bottom of every
    slide, which looks like a bug in the deck rather than a deliberate space.
    """
    tip_shapes = [
        s for s in slide.shapes
        if s.has_text_frame and TIP_RE.match(s.text_frame.text or "")
    ]
    doomed = list(tip_shapes)
    for tip in tip_shapes:
        for other in slide.shapes:
            if other in doomed:
                continue
            if not (other.has_text_frame and other.text_frame.text.strip()) \
                    and backs_onto(other, tip):
                doomed.append(other)
    for shape in doomed:
        shape._element.getparent().remove(shape._element)
    return len(tip_shapes)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


STRAY_EYEBROW = re.compile(r"^\W*READ ME FIRST\W*$", re.IGNORECASE)


def repair_eyebrows(slides):
    """Fix eyebrow labels that still say READ ME FIRST on a content slide.

    One shipped template has this slip. Left alone it puts an instruction to
    the speaker on a slide the room is looking at, so we heal it from the
    eyebrow the rest of the deck already uses rather than inventing wording.
    """
    counts = {}
    for slide in slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if t and "·" in t and len(t) < 80 and not STRAY_EYEBROW.match(t):
                counts[t] = counts.get(t, 0) + 1
    if not counts:
        return []
    canonical = max(counts, key=counts.get)
    fixed = []
    for i, slide in enumerate(slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame and STRAY_EYEBROW.match(
                    shape.text_frame.text.strip()):
                set_paragraph_text(shape.text_frame.paragraphs[0], canonical)
                fixed.append(i)
    return fixed


# --------------------------------------------------------------------------
# commands
# --------------------------------------------------------------------------
def load_template(name):
    if name not in RITUALS:
        sys.exit(f"unknown ritual {name!r}. one of: {', '.join(sorted(RITUALS))}")
    path = os.path.join(TEMPLATE_DIR, RITUALS[name])
    if not os.path.exists(path):
        sys.exit(f"template missing: {path}")
    return Presentation(path)


def cmd_list(name):
    prs = load_template(name)
    slides = list(prs.slides)
    print(f"# {name} — {len(slides)} slides in template, "
          f"{len(slides) - 1} in the delivered deck\n")
    for i, slide in enumerate(slides):
        label = "READ ME FIRST (deleted on fill)" if i == 0 else f"slide {i}"
        heads = [s.text_frame.text.strip().split("\n")[0]
                 for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        title = heads[1] if len(heads) > 1 else (heads[0] if heads else "")
        print(f"## {label}: {title}")
        if i == 0:
            print()
            continue
        tokens = []
        for _, tf in iter_text_frames(slide):
            for tok in re.findall(r"\[[^\[\]]+\]", tf.text):
                if tok not in tokens:
                    tokens.append(tok)
        for tok in tokens:
            print(f"   token  {tok}")
        for t_i, shape in enumerate(s for s in slide.shapes if s.has_table):
            tbl = shape.table
            print(f"   table  index={t_i}  {len(tbl.rows)} rows x "
                  f"{len(tbl.columns)} cols")
            for row in tbl.rows:
                print("          | " + " | ".join(
                    c.text.replace("\n", " / ")[:40] for c in row.cells))
        notes = slide.notes_slide.notes_text_frame.text.strip() \
            if slide.has_notes_slide else ""
        if notes:
            print(f"   notes  {notes[:160].replace(chr(10), ' ')}")
        print()


def cmd_fill(name, data_path, out_path, keep_readme=False, keep_tips=False):
    prs = load_template(name)
    with open(data_path) as fh:
        data = json.load(fh)

    if not keep_readme:
        drop_slide(prs, 0)

    slides = list(prs.slides)
    report = []

    for entry in data.get("slides", []):
        n = entry["n"]
        if not 1 <= n <= len(slides):
            report.append(f"! slide {n} does not exist (deck has {len(slides)})")
            continue
        slide = slides[n - 1]
        for token, value in (entry.get("replace") or {}).items():
            hits = replace_token(slide, token, value)
            if hits == 0:
                report.append(f"! slide {n}: token {token!r} not found")
        for spec in entry.get("tables") or []:
            if not isinstance(spec, dict) or "rows" not in spec:
                report.append(
                    f"! slide {n}: 'tables' must be a list of "
                    f'{{"index": 0, "rows": [[...]]}} objects, got {spec!r}')
                continue
            tables = [s for s in slide.shapes if s.has_table]
            idx = spec.get("index", 0)
            if idx >= len(tables):
                report.append(f"! slide {n}: no table at index {idx}")
                continue
            fill_table(tables[idx], spec["rows"])
        if entry.get("notes"):
            set_notes(slide, entry["notes"])

    if not keep_tips:
        for slide in slides:
            strip_tips(slide)

    for i in repair_eyebrows(slides):
        report.append(f"  slide {i}: eyebrow label said READ ME FIRST, healed")

    prs.save(out_path)
    print(f"saved {out_path} — {len(slides)} slides")
    for line in report:
        print(line)
    if any(line.startswith("!") for line in report):
        print("\nTokens that were not found usually mean the wording drifted. "
              "Run --list and copy the token exactly.")
    return 0


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--template", required=True, help=", ".join(sorted(RITUALS)))
    ap.add_argument("--list", action="store_true",
                    help="print every fillable token and table, then exit")
    ap.add_argument("--data", help="JSON content file")
    ap.add_argument("--out", help="output .pptx")
    ap.add_argument("--keep-readme", action="store_true")
    ap.add_argument("--keep-tips", action="store_true")
    args = ap.parse_args()

    if args.list:
        cmd_list(args.template)
        return
    if not (args.data and args.out):
        sys.exit("need --data and --out (or --list)")
    cmd_fill(args.template, args.data, args.out, args.keep_readme, args.keep_tips)


if __name__ == "__main__":
    main()
